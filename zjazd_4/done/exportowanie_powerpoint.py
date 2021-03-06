from pptx import Presentation

prs = Presentation()

slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

title_shape = shapes.title

body_shape = shapes.placeholders[1]

title_shape.text = "Jakis text"

tf = body_shape.text_frame
tf.tex = "Zawartosc text frame"

with open("raport.csv") as f:
    data = csv.reader(f, delimiter=",")
    for row in data:
        p = tf.add_paragraph()
        p.text = row[0]
        p.level = 1

        p = tf.add_paragraph()
        p.text = row[1]
        p.level = 2

prs.save("raport.pptx")